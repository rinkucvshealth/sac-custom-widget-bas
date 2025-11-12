"""
Convert BASIS_TEAM_PRESENTATION.md to PowerPoint (.pptx)
Install required: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re

def parse_markdown_presentation(file_path):
    """Parse the markdown presentation file"""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    slides = []
    current_slide = {'title': '', 'content': []}
    
    lines = content.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        
        # Check for slide header (## Slide X:)
        if line.startswith('## Slide'):
            if current_slide['title']:
                slides.append(current_slide)
            current_slide = {'title': '', 'content': []}
            # Extract title (text after ":")
            if ':' in line:
                current_slide['title'] = line.split(':', 1)[1].strip()
        # Check for main title (# )
        elif line.startswith('# ') and not slides:
            current_slide['title'] = line[2:].strip()
        # Check for headers (## )
        elif line.startswith('### '):
            current_slide['content'].append(('header', line[4:].strip()))
        # Check for bold text (**)
        elif line.startswith('**') and line.endswith('**'):
            current_slide['content'].append(('bold', line[2:-2].strip()))
        # Check for bullet points (-)
        elif line.startswith('- '):
            text = line[2:].strip()
            # Check for bold in bullet
            if '**' in text:
                text = text.replace('**', '')
            current_slide['content'].append(('bullet', text))
        # Check for checkmarks (✅)
        elif '✅' in line or '❌' in line:
            current_slide['content'].append(('bullet', line.strip()))
        # Regular content
        elif line and not line.startswith('---'):
            if line not in ['```typescript', '```', '**Thank you for your consideration!**']:
                current_slide['content'].append(('text', line))
        
        i += 1
    
    # Add last slide
    if current_slide['title']:
        slides.append(current_slide)
    
    return slides

def create_presentation(slides_data, output_path):
    """Create PowerPoint presentation from slides data"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides_data:
        # Use Title and Content layout
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = slide_data['title']
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Set content
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.word_wrap = True
        
        # Clear default paragraphs
        tf.clear()
        
        for item_type, text in slide_data['content']:
            if not text or text.startswith('```'):
                continue
                
            if item_type == 'header':
                p = tf.add_paragraph()
                p.text = text
                p.level = 0
                p.font.size = Pt(20)
                p.font.bold = True
                p.space_after = Pt(12)
            elif item_type == 'bold':
                p = tf.add_paragraph()
                p.text = text
                p.level = 0
                p.font.size = Pt(18)
                p.font.bold = True
                p.space_after = Pt(6)
            elif item_type == 'bullet':
                p = tf.add_paragraph()
                p.text = text
                p.level = 0
                p.font.size = Pt(16)
                p.space_after = Pt(4)
                p.bullet = True
            else:  # text
                p = tf.add_paragraph()
                p.text = text
                p.level = 0
                p.font.size = Pt(14)
                p.space_after = Pt(4)
    
    prs.save(output_path)
    print(f"✅ PowerPoint presentation created: {output_path}")
    print(f"   Total slides: {len(slides_data)}")

if __name__ == '__main__':
    import sys
    import os
    
    markdown_file = '../BASIS_TEAM_PRESENTATION.md'
    output_file = '../BASIS_TEAM_PRESENTATION.pptx'
    
    if not os.path.exists(markdown_file):
        print(f"❌ File not found: {markdown_file}")
        sys.exit(1)
    
    print(f"📖 Reading: {markdown_file}")
    slides = parse_markdown_presentation(markdown_file)
    print(f"📊 Parsed {len(slides)} slides")
    
    print(f"📝 Creating PowerPoint...")
    create_presentation(slides, output_file)
    print(f"✅ Done! You can now upload {output_file} to Google Slides")




